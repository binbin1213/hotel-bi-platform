import os
import logging
import asyncio
import json
from datetime import datetime
from typing import Dict, List, Optional, Union, Any
from sqlalchemy.orm import Session

# 条件导入，如果没有安装相关库，则跳过导入
try:
    from playwright.async_api import async_playwright
except ImportError:
    logging.warning("Playwright库未安装，PDF生成功能将不可用")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.chart.data import CategoryChartData
except ImportError:
    logging.warning("python-pptx库未安装，PPT生成功能将不可用")

from ..models.report import Report
from ..models.hotel_data import HotelData
from ..models.kpi import KPIMetric
from ..config.settings import settings
from ..utils.file_handler import upload_file_to_minio, generate_download_url

logger = logging.getLogger(__name__)

class ReportService:
    """报告生成服务，负责生成PDF和PPT格式的报告"""
    
    def __init__(self, db: Session):
        self.db = db
        self.templates_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), "templates")
        self.temp_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), "temp")
        
        # 确保临时目录存在
        os.makedirs(self.temp_dir, exist_ok=True)
        os.makedirs(self.templates_dir, exist_ok=True)
    
    async def generate_pdf_report(self, report_id: int) -> str:
        """
        使用Playwright渲染HTML生成PDF报告
        
        Args:
            report_id: 报告ID
            
        Returns:
            str: PDF文件的下载URL
        """
        logger.info(f"开始生成PDF报告，报告ID: {report_id}")
        
        try:
            # 检查Playwright是否可用
            from playwright.async_api import async_playwright
        except ImportError:
            raise RuntimeError("Playwright库未安装，无法生成PDF报告")
        
        # 获取报告数据
        report = self.db.query(Report).filter(Report.id == report_id).first()
        if not report:
            raise ValueError(f"报告不存在，ID: {report_id}")
        
        # 获取报告相关数据
        report_data = self._prepare_report_data(report)
        
        # 生成HTML内容
        html_content = self._render_html_template(report_data)
        
        # 临时HTML文件路径
        html_path = os.path.join(self.temp_dir, f"report_{report_id}_{datetime.now().strftime('%Y%m%d%H%M%S')}.html")
        
        # 保存HTML文件
        with open(html_path, "w", encoding="utf-8") as f:
            f.write(html_content)
        
        # PDF文件路径
        pdf_path = html_path.replace(".html", ".pdf")
        
        try:
            # 使用Playwright渲染HTML为PDF
            async with async_playwright() as p:
                browser = await p.chromium.launch()
                page = await browser.new_page()
                await page.goto(f"file://{html_path}")
                await page.pdf(path=pdf_path, format="A4")
                await browser.close()
            
            # 上传PDF到MinIO
            file_key = f"reports/pdf/{os.path.basename(pdf_path)}"
            upload_file_to_minio(pdf_path, file_key)
            
            # 生成下载URL
            download_url = generate_download_url(file_key)
            
            # 更新报告文件路径
            if report.file_paths:
                file_paths = json.loads(report.file_paths)
            else:
                file_paths = {}
            
            file_paths["pdf"] = file_key
            report.file_paths = json.dumps(file_paths)
            self.db.commit()
            
            # 删除临时文件
            os.remove(html_path)
            os.remove(pdf_path)
            
            logger.info(f"PDF报告生成成功，报告ID: {report_id}")
            return download_url
            
        except Exception as e:
            logger.error(f"生成PDF报告失败: {str(e)}")
            # 清理临时文件
            if os.path.exists(html_path):
                os.remove(html_path)
            if os.path.exists(pdf_path):
                os.remove(pdf_path)
            raise RuntimeError(f"生成PDF报告失败: {str(e)}")
    
    def generate_ppt_report(self, report_id: int) -> str:
        """
        使用python-pptx生成PPT报告
        
        Args:
            report_id: 报告ID
            
        Returns:
            str: PPT文件的下载URL
        """
        logger.info(f"开始生成PPT报告，报告ID: {report_id}")
        
        try:
            # 检查python-pptx是否可用
            from pptx import Presentation
        except ImportError:
            raise RuntimeError("python-pptx库未安装，无法生成PPT报告")
        
        # 获取报告数据
        report = self.db.query(Report).filter(Report.id == report_id).first()
        if not report:
            raise ValueError(f"报告不存在，ID: {report_id}")
        
        # 获取报告相关数据
        report_data = self._prepare_report_data(report)
        
        # PPT模板路径
        template_path = os.path.join(self.templates_dir, "report_template.pptx")
        if not os.path.exists(template_path):
            # 如果模板不存在，创建一个基础模板
            template_path = self._create_base_template()
        
        # 创建PPT演示文稿
        prs = Presentation(template_path)
        
        # 添加标题幻灯片
        self._add_title_slide(prs, report_data)
        
        # 添加摘要幻灯片
        self._add_summary_slide(prs, report_data)
        
        # 添加KPI指标幻灯片
        self._add_kpi_slide(prs, report_data)
        
        # 添加趋势分析幻灯片
        self._add_trend_slide(prs, report_data)
        
        # 添加AI分析幻灯片
        self._add_ai_insights_slide(prs, report_data)
        
        # 添加建议幻灯片
        self._add_recommendations_slide(prs, report_data)
        
        # PPT文件路径
        ppt_path = os.path.join(self.temp_dir, f"report_{report_id}_{datetime.now().strftime('%Y%m%d%H%M%S')}.pptx")
        
        # 保存PPT文件
        prs.save(ppt_path)
        
        try:
            # 上传PPT到MinIO
            file_key = f"reports/ppt/{os.path.basename(ppt_path)}"
            upload_file_to_minio(ppt_path, file_key)
            
            # 生成下载URL
            download_url = generate_download_url(file_key)
            
            # 更新报告文件路径
            if report.file_paths:
                file_paths = json.loads(report.file_paths)
            else:
                file_paths = {}
            
            file_paths["ppt"] = file_key
            report.file_paths = json.dumps(file_paths)
            self.db.commit()
            
            # 删除临时文件
            os.remove(ppt_path)
            
            logger.info(f"PPT报告生成成功，报告ID: {report_id}")
            return download_url
            
        except Exception as e:
            logger.error(f"生成PPT报告失败: {str(e)}")
            # 清理临时文件
            if os.path.exists(ppt_path):
                os.remove(ppt_path)
            raise RuntimeError(f"生成PPT报告失败: {str(e)}")
    
    def _prepare_report_data(self, report: Report) -> Dict[str, Any]:
        """准备报告所需的数据"""
        report_data = {
            "id": report.id,
            "title": report.title,
            "report_type": report.report_type,
            "created_at": report.created_at.strftime("%Y-%m-%d %H:%M:%S"),
            "ai_insights": report.ai_insights,
        }
        
        # 解析报告内容数据
        if report.content_data:
            content_data = json.loads(report.content_data)
            report_data.update(content_data)
        else:
            # 如果报告没有内容数据，获取基础数据
            hotel_data = self._get_hotel_data_for_report(report)
            kpi_metrics = self._get_kpi_metrics_for_report(report)
            
            report_data["hotel_data"] = hotel_data
            report_data["kpi_metrics"] = kpi_metrics
        
        return report_data
    
    def _get_hotel_data_for_report(self, report: Report) -> List[Dict[str, Any]]:
        """获取报告相关的酒店数据"""
        # 这里可以根据报告类型和内容获取相关的酒店数据
        # 示例实现，实际应根据业务需求调整
        hotel_data_list = []
        
        # 假设报告内容中包含了酒店ID或日期范围等信息
        if report.content_data:
            content = json.loads(report.content_data)
            hotel_ids = content.get("hotel_ids", [])
            date_from = content.get("date_from")
            date_to = content.get("date_to")
            
            query = self.db.query(HotelData)
            
            if hotel_ids:
                query = query.filter(HotelData.id.in_(hotel_ids))
            
            if date_from and date_to:
                query = query.filter(HotelData.date_recorded.between(date_from, date_to))
            
            hotel_data_records = query.all()
            
            for record in hotel_data_records:
                hotel_data_list.append(record.to_dict())
        
        return hotel_data_list
    
    def _get_kpi_metrics_for_report(self, report: Report) -> List[Dict[str, Any]]:
        """获取报告相关的KPI指标"""
        # 示例实现，实际应根据业务需求调整
        kpi_metrics_list = []
        
        if report.content_data:
            content = json.loads(report.content_data)
            hotel_ids = content.get("hotel_ids", [])
            metric_types = content.get("metric_types", [])
            date_from = content.get("date_from")
            date_to = content.get("date_to")
            
            query = self.db.query(KPIMetric)
            
            if hotel_ids:
                query = query.filter(KPIMetric.hotel_id.in_(hotel_ids))
            
            if metric_types:
                query = query.filter(KPIMetric.metric_type.in_(metric_types))
            
            if date_from and date_to:
                query = query.filter(KPIMetric.period_start >= date_from)
                query = query.filter(KPIMetric.period_end <= date_to)
            
            kpi_records = query.all()
            
            for record in kpi_records:
                kpi_metrics_list.append(record.to_dict())
        
        return kpi_metrics_list
    
    def _render_html_template(self, report_data: Dict[str, Any]) -> str:
        """渲染HTML模板"""
        # 这里可以使用Jinja2等模板引擎渲染HTML
        # 简单示例，实际应使用模板引擎
        html = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <meta charset="UTF-8">
            <title>{report_data['title']}</title>
            <style>
                body {{ font-family: Arial, sans-serif; margin: 0; padding: 20px; }}
                .header {{ text-align: center; margin-bottom: 30px; }}
                .section {{ margin-bottom: 20px; }}
                .section-title {{ color: #333; border-bottom: 1px solid #ddd; padding-bottom: 5px; }}
                .kpi-grid {{ display: grid; grid-template-columns: repeat(3, 1fr); gap: 15px; }}
                .kpi-card {{ background: #f5f5f5; padding: 15px; border-radius: 5px; }}
                .kpi-value {{ font-size: 24px; font-weight: bold; color: #0066cc; }}
                .insights {{ background: #f0f7ff; padding: 15px; border-left: 4px solid #0066cc; }}
            </style>
        </head>
        <body>
            <div class="header">
                <h1>{report_data['title']}</h1>
                <p>生成时间: {report_data['created_at']}</p>
            </div>
            
            <div class="section">
                <h2 class="section-title">摘要</h2>
                <p>本报告提供了酒店运营数据的分析结果，包括关键绩效指标、趋势分析和AI洞察。</p>
            </div>
            
            <div class="section">
                <h2 class="section-title">关键绩效指标</h2>
                <div class="kpi-grid">
                    <!-- 这里可以根据实际KPI数据动态生成 -->
                    <div class="kpi-card">
                        <h3>入住率</h3>
                        <div class="kpi-value">85%</div>
                    </div>
                    <div class="kpi-card">
                        <h3>平均房价</h3>
                        <div class="kpi-value">¥688</div>
                    </div>
                    <div class="kpi-card">
                        <h3>RevPAR</h3>
                        <div class="kpi-value">¥584</div>
                    </div>
                </div>
            </div>
            
            <div class="section">
                <h2 class="section-title">AI分析洞察</h2>
                <div class="insights">
                    <p>{report_data.get('ai_insights', '暂无AI分析')}</p>
                </div>
            </div>
            
            <div class="section">
                <h2 class="section-title">建议</h2>
                <ul>
                    <li>根据入住率数据，建议在淡季增加促销活动</li>
                    <li>平均房价高于行业平均水平，具有良好的市场竞争力</li>
                    <li>RevPAR指标表现优秀，建议维持当前定价策略</li>
                </ul>
            </div>
        </body>
        </html>
        """
        return html
    
    def _create_base_template(self) -> str:
        """创建基础PPT模板"""
        prs = Presentation()
        
        # 保存基础模板
        template_path = os.path.join(self.templates_dir, "report_template.pptx")
        os.makedirs(os.path.dirname(template_path), exist_ok=True)
        prs.save(template_path)
        
        return template_path
    
    def _add_title_slide(self, prs: Presentation, report_data: Dict[str, Any]) -> None:
        """添加标题幻灯片"""
        title_slide_layout = prs.slide_layouts[0]  # 标题幻灯片布局
        slide = prs.slides.add_slide(title_slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = report_data["title"]
        subtitle.text = f"生成时间: {report_data['created_at']}"
    
    def _add_summary_slide(self, prs: Presentation, report_data: Dict[str, Any]) -> None:
        """添加摘要幻灯片"""
        bullet_slide_layout = prs.slide_layouts[1]  # 标题和内容布局
        slide = prs.slides.add_slide(bullet_slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "报告摘要"
        
        tf = content.text_frame
        tf.text = "本报告包含以下内容:"
        
        p = tf.add_paragraph()
        p.text = "• 关键绩效指标分析"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• 趋势分析"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• AI智能洞察"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• 建议和行动计划"
        p.level = 1
    
    def _add_kpi_slide(self, prs: Presentation, report_data: Dict[str, Any]) -> None:
        """添加KPI指标幻灯片"""
        slide_layout = prs.slide_layouts[2]  # 标题和内容布局
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "关键绩效指标"
        
        # 创建KPI表格
        rows, cols = 4, 3
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(3)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # 设置表头
        table.cell(0, 0).text = "指标"
        table.cell(0, 1).text = "当前值"
        table.cell(0, 2).text = "同比变化"
        
        # 填充KPI数据（示例数据）
        kpi_data = [
            ("入住率", "85%", "+5%"),
            ("平均房价", "¥688", "+3%"),
            ("RevPAR", "¥584", "+8%")
        ]
        
        for i, (metric, value, change) in enumerate(kpi_data, start=1):
            table.cell(i, 0).text = metric
            table.cell(i, 1).text = value
            table.cell(i, 2).text = change
    
    def _add_trend_slide(self, prs: Presentation, report_data: Dict[str, Any]) -> None:
        """添加趋势分析幻灯片"""
        slide_layout = prs.slide_layouts[5]  # 标题和内容布局
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "趋势分析"
        
        # 在实际应用中，这里应该添加图表
        # 简单示例：添加文本框代替图表
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(5)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.text = "这里将显示趋势图表，展示关键指标的历史变化趋势。"
        
        p = tf.add_paragraph()
        p.text = "• 入住率呈现稳定上升趋势"
        
        p = tf.add_paragraph()
        p.text = "• 平均房价在节假日期间有明显波动"
        
        p = tf.add_paragraph()
        p.text = "• RevPAR同比增长显著"
    
    def _add_ai_insights_slide(self, prs: Presentation, report_data: Dict[str, Any]) -> None:
        """添加AI分析幻灯片"""
        slide_layout = prs.slide_layouts[1]  # 标题和内容布局
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "AI智能洞察"
        
        tf = content.text_frame
        tf.text = report_data.get("ai_insights", "暂无AI分析")
    
    def _add_recommendations_slide(self, prs: Presentation, report_data: Dict[str, Any]) -> None:
        """添加建议幻灯片"""
        slide_layout = prs.slide_layouts[1]  # 标题和内容布局
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "建议和行动计划"
        
        tf = content.text_frame
        tf.text = "基于数据分析和AI洞察，我们提出以下建议:"
        
        # 示例建议，实际应基于AI分析结果
        p = tf.add_paragraph()
        p.text = "• 根据入住率数据，建议在淡季增加促销活动"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• 平均房价高于行业平均水平，具有良好的市场竞争力"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• RevPAR指标表现优秀，建议维持当前定价策略"
        p.level = 1 